import asyncio
import base64
import io
import logging
import os
import threading
import urllib.parse
import uuid
from datetime import datetime, timedelta
from http.server import HTTPServer, BaseHTTPRequestHandler

import aiohttp
import asyncpg
import openpyxl
from docx import Document
from docx.shared import Pt, RGBColor
from anthropic import AsyncAnthropic
from aiogram import Bot, Dispatcher, F
from aiogram.client.session.aiohttp import AiohttpSession
from aiogram.filters import CommandStart, Command
from aiogram.types import (
    Message, BufferedInputFile,
    LabeledPrice, PreCheckoutQuery,
    InlineKeyboardMarkup, InlineKeyboardButton,
    ReplyKeyboardMarkup, KeyboardButton
)
from pptx import Presentation

# ─── Настройки ────────────────────────────────────────────────────────────────
TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN", "YOUR_TELEGRAM_BOT_TOKEN")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "YOUR_ANTHROPIC_API_KEY")
PROXY_URL = os.getenv("PROXY_URL", "")
WOLFRAM_API_KEY = os.getenv("WOLFRAM_API_KEY", "")
YUKASSA_SHOP_ID = os.getenv("YUKASSA_SHOP_ID", "")
YUKASSA_SECRET_KEY = os.getenv("YUKASSA_SECRET_KEY", "")
TOGETHER_API_KEY = os.getenv("TOGETHER_API_KEY", "")
TAVILY_API_KEY = os.getenv("TAVILY_API_KEY", "")

# Цена подписки
SUBSCRIPTION_PRICE_RUB = 299
SUBSCRIPTION_PRICE_STARS = 250
SUBSCRIPTION_DAYS = 30
FREE_MESSAGES_PER_DAY = 10
TRIAL_DAYS = 3
DAILY_BONUS_MESSAGES = 3
MAX_BONUS_MESSAGES = 15
MAX_HISTORY = 20
ADMIN_IDS = [int(x) for x in os.getenv("ADMIN_IDS", "0").split(",") if x and x != "0"]

# ─── Инициализация ────────────────────────────────────────────────────────────
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

if PROXY_URL:
    session = AiohttpSession(proxy=PROXY_URL)
    bot = Bot(token=TELEGRAM_TOKEN, session=session)
else:
    bot = Bot(token=TELEGRAM_TOKEN)

dp = Dispatcher()
anthropic = AsyncAnthropic(api_key=ANTHROPIC_API_KEY)

# ─── ЮKassa платежи ──────────────────────────────────────────────────────────
async def create_yukassa_payment(user_id: int, username: str) -> dict | None:
    """Создаём платёж в ЮKassa и возвращаем ссылку для оплаты."""
    if not YUKASSA_SHOP_ID or not YUKASSA_SECRET_KEY:
        return None
    try:
        payment_id = str(uuid.uuid4())
        payload = {
            "amount": {
                "value": f"{SUBSCRIPTION_PRICE_RUB}.00",
                "currency": "RUB"
            },
            "confirmation": {
                "type": "redirect",
                "return_url": f"https://t.me/{(await bot.get_me()).username}"
            },
            "capture": True,
            "description": f"Подписка NeuroBot на 30 дней — @{username}",
            "metadata": {
                "user_id": str(user_id)
            }
        }
        auth = aiohttp.BasicAuth(YUKASSA_SHOP_ID, YUKASSA_SECRET_KEY)
        headers = {
            "Idempotence-Key": payment_id,
            "Content-Type": "application/json"
        }
        async with aiohttp.ClientSession() as s:
            async with s.post(
                "https://api.yookassa.ru/v3/payments",
                json=payload,
                auth=auth,
                headers=headers,
                timeout=aiohttp.ClientTimeout(total=15)
            ) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    return {
                        "payment_id": data["id"],
                        "url": data["confirmation"]["confirmation_url"]
                    }
                else:
                    body = await resp.text()
                    logger.error(f"ЮKassa ошибка: {resp.status} {body}")
    except Exception as e:
        logger.error(f"ЮKassa exception: {e}")
    return None


async def check_yukassa_payment(payment_id: str) -> bool:
    """Проверяем статус платежа."""
    if not YUKASSA_SHOP_ID or not YUKASSA_SECRET_KEY:
        return False
    try:
        auth = aiohttp.BasicAuth(YUKASSA_SHOP_ID, YUKASSA_SECRET_KEY)
        async with aiohttp.ClientSession() as s:
            async with s.get(
                f"https://api.yookassa.ru/v3/payments/{payment_id}",
                auth=auth,
                timeout=aiohttp.ClientTimeout(total=15)
            ) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    return data.get("status") == "succeeded"
    except Exception as e:
        logger.error(f"ЮKassa check error: {e}")
    return False


# Хранилище ожидающих платежей {payment_id: user_id}
pending_payments: dict[str, int] = {}


def save_payment(user_id: int, username: str, payment_id: str):
    """Сохраняем успешный платёж в БД."""
    conn = sqlite3.connect("bot.db")
    c = conn.cursor()
    now = datetime.now().isoformat()
    c.execute(
        "INSERT INTO payments (user_id, username, amount, payment_id, created_at) VALUES (?, ?, ?, ?, ?)",
        (user_id, username, SUBSCRIPTION_PRICE_RUB, payment_id, now)
    )
    conn.commit()
    conn.close()


async def poll_payment(payment_id: str, user_id: int):
    """Проверяем платёж каждые 30 секунд в течение 30 минут."""
    for _ in range(60):
        await asyncio.sleep(30)
        paid = await check_yukassa_payment(payment_id)
        if paid:
            await activate_subscription(user_id, from_payment=True)
            until = await get_subscription_until(user_id)
            username = (await bot.get_chat(user_id)).username or str(user_id)
            await save_payment(user_id, username, payment_id)
            try:
                await bot.send_message(
                    user_id,
                    f"🎉 Оплата прошла успешно!\n\n"
                    f"✅ Подписка активна до {until}\n\n"
                    f"Теперь у тебя безлимитный доступ ко всем функциям!"
                )
                for admin_id in ADMIN_IDS:
                    try:
                        await bot.send_message(
                            admin_id,
                            f"💰 Новая оплата!\n\n"
                            f"👤 Пользователь: @{username}\n"
                            f"💵 Сумма: {SUBSCRIPTION_PRICE_RUB}₽\n"
                            f"📅 Подписка до: {until}"
                        )
                    except Exception:
                        pass
            except Exception as e:
                logger.error(f"Ошибка уведомления об оплате: {e}")
            pending_payments.pop(payment_id, None)
            return
    pending_payments.pop(payment_id, None)
class HealthHandler(BaseHTTPRequestHandler):
    def do_GET(self):
        self.send_response(200)
        self.end_headers()
        self.wfile.write(b"OK")
    def log_message(self, *args):
        pass  # Отключаем логи

def run_health_server():
    server = HTTPServer(("0.0.0.0", 8080), HealthHandler)
    server.serve_forever()


SYSTEM_PROMPT_BASE = """Ты умный AI-ассистент в Telegram с расширенными возможностями.

ВАЖНО: Ты УМЕЕШЬ генерировать изображения через встроенный инструмент! 
Когда тебя спрашивают умеешь ли ты рисовать/генерировать картинки — отвечай ДА.

ВАЖНО ПРО ФОРМАТИРОВАНИЕ: Ты работаешь в Telegram. НЕ используй markdown разметку в обычных ответах — никаких звёздочек (**текст**), решёток (# заголовок), обратных кавычек (`код`). Пиши обычным текстом. Структурируй ответы через переносы строк и эмодзи если нужно.

Твои возможности:
- Отвечать на вопросы
- Генерировать изображения по описанию (команда: "нарисуй...")  
- Распознавать и анализировать картинки
- Создавать Excel таблицы
- Создавать презентации PowerPoint
- Создавать Word документы
- Точно решать математические задачи

При решении математических и геометрических задач:
1. Всегда решай строго пошагово
2. Записывай все формулы которые используешь
3. Подставляй числа явно на каждом шаге
4. В конце проверяй ответ подстановкой
5. Если задача геометрическая — сначала опиши фигуру и все известные элементы

Если пользователь пишет на русском — отвечай на русском.
если на английском — на английском.

Когда пользователь просит создать Excel таблицу — отвечай ТОЛЬКО в формате:
EXCEL_TABLE:
Заголовок1|Заголовок2|Заголовок3
Данные1|Данные2|Данные3

Когда пользователь просит создать презентацию — отвечай ТОЛЬКО в формате:
PRESENTATION:
TITLE:Название презентации
SLIDE:Заголовок слайда|Текст содержимого слайда
SLIDE:Заголовок 2|Текст 2

Когда пользователь просит создать Word документ, реферат, эссе, доклад, договор, заявление или любой другой текстовый документ — создай полноценный текст используя markdown разметку:
# для заголовков первого уровня
## для заголовков второго уровня
### для заголовков третьего уровня
- для списков
**текст** для жирного текста
Пиши полный развёрнутый текст документа.
"""

STYLE_PROMPTS = {
    "formal": "СТИЛЬ ОБЩЕНИЯ: Отвечай официально и вежливо. Используй грамотный литературный язык. Обращайся на «вы».",
    "friend": "СТИЛЬ ОБЩЕНИЯ: Общайся как близкий друг — неформально, живо, с юмором и эмодзи. Говори просто, можешь шутить, используй разговорный язык. Обращайся на «ты».",
    "short": "СТИЛЬ ОБЩЕНИЯ: Отвечай максимально кратко и по делу. Только суть, никакой воды. Минимум слов — максимум смысла.",
}

STYLE_NAMES = {
    "formal": "🎩 Формальный",
    "friend": "😊 Как с другом",
    "short": "⚡ Краткий",
}


async def get_system_prompt(user_id: int) -> str:
    style = await get_user_style(user_id)
    style_text = STYLE_PROMPTS.get(style, STYLE_PROMPTS["friend"])
    return SYSTEM_PROMPT_BASE + "\n" + style_text


DATABASE_URL = os.getenv("DATABASE_URL", "")
db_pool = None


async def get_db():
    return db_pool


# ─── База данных ──────────────────────────────────────────────────────────────
async def init_db():
    global db_pool
    db_pool = await asyncpg.create_pool(DATABASE_URL, ssl="require")

    async with db_pool.acquire() as conn:
        await conn.execute("""
            CREATE TABLE IF NOT EXISTS users (
                user_id BIGINT PRIMARY KEY,
                username TEXT,
                created_at TEXT,
                subscription_until TEXT,
                messages_today INTEGER DEFAULT 0,
                last_message_date TEXT,
                referred_by BIGINT DEFAULT NULL,
                referral_count INTEGER DEFAULT 0,
                notified_expiry INTEGER DEFAULT 0,
                trial_used INTEGER DEFAULT 0,
                bonus_messages INTEGER DEFAULT 0,
                last_bonus_date TEXT,
                speech_style TEXT DEFAULT 'friend'
            )
        """)
        await conn.execute("""
            CREATE TABLE IF NOT EXISTS message_history (
                id SERIAL PRIMARY KEY,
                user_id BIGINT,
                role TEXT,
                content TEXT,
                created_at TEXT
            )
        """)
        await conn.execute("""
            CREATE TABLE IF NOT EXISTS payments (
                id SERIAL PRIMARY KEY,
                user_id BIGINT,
                username TEXT,
                amount INTEGER,
                payment_id TEXT,
                created_at TEXT
            )
        """)


async def get_user(user_id: int) -> dict | None:
    async with db_pool.acquire() as conn:
        row = await conn.fetchrow("SELECT * FROM users WHERE user_id = $1", user_id)
        if row:
            return dict(row)
        return None


async def create_user(user_id: int, username: str, referred_by: int = None):
    async with db_pool.acquire() as conn:
        now = datetime.now().isoformat()
        today = datetime.now().date().isoformat()
        result = await conn.execute("""
            INSERT INTO users (user_id, username, created_at, messages_today, last_message_date, referred_by)
            VALUES ($1, $2, $3, 0, $4, $5)
            ON CONFLICT (user_id) DO NOTHING
        """, user_id, username, now, today, referred_by)

        if result == "INSERT 0 1" and referred_by:
            await conn.execute(
                "UPDATE users SET referral_count = referral_count + 1 WHERE user_id = $1",
                referred_by
            )
            row = await conn.fetchrow("SELECT referral_count, username FROM users WHERE user_id = $1", referred_by)
            if row:
                ref_count = row["referral_count"]
                ref_username = row["username"] or str(referred_by)

                # Уведомление админу когда кто-то достигает 10 рефералов
                if ref_count == 10:
                    for admin_id in ADMIN_IDS:
                        try:
                            await bot.send_message(
                                admin_id,
                                f"🏆 Победитель конкурса!\n\n"
                                f"👤 @{ref_username} привёл 10 друзей!\n\n"
                                f"Не забудь выдать приз:\n"
                                f"• Месяц подписки\n"
                                f"• Ранний доступ к голосовым"
                            )
                        except Exception:
                            pass

                # Каждые 5 рефералов — бесплатный месяц
                if ref_count % 5 == 0:
                    await activate_subscription(referred_by)
                    return True
    return False


async def get_referral_count(user_id: int) -> int:
    async with db_pool.acquire() as conn:
        row = await conn.fetchrow("SELECT referral_count FROM users WHERE user_id = $1", user_id)
        return row["referral_count"] if row else 0


async def get_user_style(user_id: int) -> str:
    async with db_pool.acquire() as conn:
        row = await conn.fetchrow("SELECT speech_style FROM users WHERE user_id = $1", user_id)
        return row["speech_style"] if row and row["speech_style"] else "friend"


async def set_user_style(user_id: int, style: str):
    async with db_pool.acquire() as conn:
        await conn.execute("UPDATE users SET speech_style = $1 WHERE user_id = $2", style, user_id)


async def activate_trial(user_id: int) -> bool:
    async with db_pool.acquire() as conn:
        trial_until = (datetime.now() + timedelta(days=TRIAL_DAYS)).isoformat()
        result = await conn.execute("""
            UPDATE users SET subscription_until = $1, trial_used = 1
            WHERE user_id = $2 AND trial_used = 0
        """, trial_until, user_id)
        return result == "UPDATE 1"


async def is_subscribed(user_id: int) -> bool:
    user = await get_user(user_id)
    if not user or not user.get("subscription_until"):
        return False
    until = datetime.fromisoformat(user["subscription_until"])
    return until > datetime.now()


async def get_subscription_until(user_id: int) -> str | None:
    user = await get_user(user_id)
    if not user or not user.get("subscription_until"):
        return None
    until = datetime.fromisoformat(user["subscription_until"])
    if until > datetime.now():
        return until.strftime("%d.%m.%Y")
    return None


async def activate_subscription(user_id: int, from_payment: bool = False):
    async with db_pool.acquire() as conn:
        user = await get_user(user_id)
        if not from_payment and user and user.get("subscription_until"):
            current = datetime.fromisoformat(user["subscription_until"])
            if current > datetime.now():
                new_until = current + timedelta(days=SUBSCRIPTION_DAYS)
            else:
                new_until = datetime.now() + timedelta(days=SUBSCRIPTION_DAYS)
        else:
            # При оплате всегда считаем от сегодня
            new_until = datetime.now() + timedelta(days=SUBSCRIPTION_DAYS)
        await conn.execute(
            "UPDATE users SET subscription_until = $1 WHERE user_id = $2",
            new_until.isoformat(), user_id
        )


async def can_send_message(user_id: int) -> tuple[bool, int]:
    if await is_subscribed(user_id):
        return True, -1

    user = await get_user(user_id)
    if not user:
        return False, 0

    today = datetime.now().date().isoformat()
    bonus = user.get("bonus_messages", 0) or 0
    total_limit = FREE_MESSAGES_PER_DAY + bonus

    if user.get("last_message_date") != today:
        async with db_pool.acquire() as conn:
            await conn.execute(
                "UPDATE users SET messages_today = 0, last_message_date = $1 WHERE user_id = $2",
                today, user_id
            )
        return True, total_limit

    remaining = total_limit - (user.get("messages_today") or 0)
    return remaining > 0, max(0, remaining)


async def increment_message_count(user_id: int):
    today = datetime.now().date().isoformat()
    async with db_pool.acquire() as conn:
        await conn.execute(
            "UPDATE users SET messages_today = messages_today + 1, last_message_date = $1 WHERE user_id = $2",
            today, user_id
        )


async def claim_daily_bonus(user_id: int) -> bool:
    async with db_pool.acquire() as conn:
        today = datetime.now().date().isoformat()
        row = await conn.fetchrow(
            "SELECT last_bonus_date, bonus_messages FROM users WHERE user_id = $1", user_id
        )
        if not row or row["last_bonus_date"] == today:
            return False
        current_bonus = row["bonus_messages"] or 0
        if current_bonus >= MAX_BONUS_MESSAGES:
            await conn.execute(
                "UPDATE users SET last_bonus_date = $1 WHERE user_id = $2", today, user_id
            )
            return False
        new_bonus = min(current_bonus + DAILY_BONUS_MESSAGES, MAX_BONUS_MESSAGES)
        await conn.execute(
            "UPDATE users SET bonus_messages = $1, last_bonus_date = $2 WHERE user_id = $3",
            new_bonus, today, user_id
        )
        return True


async def get_users_expiring_soon() -> list:
    async with db_pool.acquire() as conn:
        now = datetime.now()
        in_3_days = (now + timedelta(days=3)).isoformat()
        tomorrow = (now + timedelta(days=1)).isoformat()
        rows = await conn.fetch("""
            SELECT user_id, username, subscription_until FROM users
            WHERE subscription_until BETWEEN $1 AND $2
            AND notified_expiry = 0
        """, tomorrow, in_3_days)
        return [(r["user_id"], r["username"], r["subscription_until"]) for r in rows]


async def mark_notified(user_id: int):
    async with db_pool.acquire() as conn:
        await conn.execute("UPDATE users SET notified_expiry = 1 WHERE user_id = $1", user_id)


async def reset_notified_expired():
    async with db_pool.acquire() as conn:
        now = datetime.now().isoformat()
        await conn.execute(
            "UPDATE users SET notified_expiry = 0 WHERE subscription_until < $1", now
        )


async def is_trial_used(user_id: int) -> bool:
    user = await get_user(user_id)
    if not user:
        return False
    return bool(user.get("trial_used", 0))


async def save_payment(user_id: int, username: str, payment_id: str):
    async with db_pool.acquire() as conn:
        now = datetime.now().isoformat()
        await conn.execute(
            "INSERT INTO payments (user_id, username, amount, payment_id, created_at) VALUES ($1, $2, $3, $4, $5)",
            user_id, username, SUBSCRIPTION_PRICE_RUB, payment_id, now
        )


# ─── История сообщений (персистентная) ───────────────────────────────────────
async def load_history(user_id: int) -> list[dict]:
    async with db_pool.acquire() as conn:
        rows = await conn.fetch(
            "SELECT role, content FROM message_history WHERE user_id = $1 ORDER BY id DESC LIMIT $2",
            user_id, MAX_HISTORY
        )
        return [{"role": r["role"], "content": r["content"]} for r in reversed(rows)]


async def save_message(user_id: int, role: str, content: str):
    async with db_pool.acquire() as conn:
        now = datetime.now().isoformat()
        await conn.execute(
            "INSERT INTO message_history (user_id, role, content, created_at) VALUES ($1, $2, $3, $4)",
            user_id, role, content if isinstance(content, str) else str(content), now
        )
        await conn.execute("""
            DELETE FROM message_history WHERE id IN (
                SELECT id FROM message_history WHERE user_id = $1
                ORDER BY id DESC OFFSET $2
            )
        """, user_id, MAX_HISTORY * 2)


async def clear_history(user_id: int):
    async with db_pool.acquire() as conn:
        await conn.execute("DELETE FROM message_history WHERE user_id = $1", user_id)


def bottom_menu() -> ReplyKeyboardMarkup:
    return ReplyKeyboardMarkup(
        keyboard=[[KeyboardButton(text="📋 Меню")]],
        resize_keyboard=True,
        persistent=True
    )


def main_keyboard() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="💳 Подписка", callback_data="menu_subscribe"),
            InlineKeyboardButton(text="📊 Статус", callback_data="menu_status"),
        ],
        [
            InlineKeyboardButton(text="👥 Пригласить друга", callback_data="menu_referral"),
            InlineKeyboardButton(text="🗑️ Очистить историю", callback_data="menu_clear"),
        ],
        [
            InlineKeyboardButton(text="🎙️ Настройки речи", callback_data="menu_speech"),
            InlineKeyboardButton(text="❓ Помощь", callback_data="menu_help"),
        ]
    ])


def speech_style_keyboard() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎩 Формальный", callback_data="style_formal")],
        [InlineKeyboardButton(text="😊 Как с другом", callback_data="style_friend")],
        [InlineKeyboardButton(text="⚡ Краткий", callback_data="style_short")],
    ])


def subscription_keyboard() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[[
        InlineKeyboardButton(
            text="💳 Купить подписку — 299₽ / месяц",
            callback_data="buy_subscription"
        )
    ]])


# ─── Вспомогательные функции ──────────────────────────────────────────────────
def is_image_request(text: str) -> bool:
    keywords = ["нарисуй", "сгенерируй картинку", "создай изображение",
                "generate image", "draw", "создай картинку", "нарисуй мне",
                "картинку с", "изображение с"]
    return any(kw in text.lower() for kw in keywords)


def is_excel_request(text: str) -> bool:
    keywords = ["excel", "таблицу", "xlsx", "создай таблицу", "сделай таблицу"]
    return any(kw in text.lower() for kw in keywords)


def is_presentation_request(text: str) -> bool:
    keywords = ["презентацию", "powerpoint", "pptx", "слайды", "создай презентацию"]
    return any(kw in text.lower() for kw in keywords)


def is_word_request(text: str) -> bool:
    keywords = ["word", "docx", "документ", "создай документ", "сделай документ",
                "напиши документ", "оформи документ", "реферат", "доклад", "эссе",
                "сочинение", "договор", "заявление", "резюме"]
    return any(kw in text.lower() for kw in keywords)


def create_word(text: str) -> bytes | None:
    """Создаём Word документ из ответа Claude."""
    try:
        lines = []
        in_doc = False
        title = "Документ"

        for line in text.split("\n"):
            if "WORD_DOCUMENT:" in line:
                in_doc = True
                continue
            if in_doc:
                lines.append(line)

        if not lines:
            # Если нет маркера — используем весь текст
            lines = text.split("\n")

        doc = Document()

        # Стиль документа
        style = doc.styles["Normal"]
        style.font.name = "Times New Roman"
        style.font.size = Pt(12)

        for line in lines:
            line = line.strip()
            if not line:
                doc.add_paragraph("")
                continue

            # Заголовок первого уровня
            if line.startswith("# "):
                heading = doc.add_heading(line[2:], level=1)
                title = line[2:]
            # Заголовок второго уровня
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            # Заголовок третьего уровня
            elif line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            # Жирный текст
            elif line.startswith("**") and line.endswith("**"):
                p = doc.add_paragraph()
                run = p.add_run(line[2:-2])
                run.bold = True
            # Список
            elif line.startswith("- ") or line.startswith("• "):
                doc.add_paragraph(line[2:], style="List Bullet")
            elif line.startswith("1. ") or (len(line) > 2 and line[0].isdigit() and line[1] == "."):
                doc.add_paragraph(line[3:] if line.startswith("1. ") else line[2:], style="List Number")
            else:
                doc.add_paragraph(line)

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except Exception as e:
        logger.error(f"Ошибка создания Word: {e}")
        return None


def is_math_request(text: str) -> bool:
    """Определяем запросы где нужно точное вычисление, а не объяснение."""
    text_lower = text.lower()

    # Исключаем теоретические вопросы
    theory_keywords = ["что такое", "что значит", "объясни", "расскажи", "как работает",
                       "зачем нужен", "для чего", "что означает", "определение"]
    if any(kw in text_lower for kw in theory_keywords):
        return False

    # Только конкретные вычисления
    calc_keywords = ["вычисли", "посчитай", "сколько будет", "реши уравнение",
                     "решить уравнение", "найди корни", "calculate", "solve",
                     "найди значение", "чему равно", "упрости", "разложи на множители"]
    if any(kw in text_lower for kw in calc_keywords):
        return True

    # Математические выражения с числами
    import re
    if re.search(r"\d+\s*[\+\-\*\/\^]\s*\d+", text):
        return True
    if re.search(r"(sin|cos|tan|log|sqrt|∫)\s*[\(\d]", text_lower):
        return True
    if re.search(r"x\^?\d*\s*[\+\-\=]", text_lower):
        return True

    return False


async def wolfram_calculate(query: str) -> str | None:
    """Точное вычисление через WolframAlpha."""
    if not WOLFRAM_API_KEY:
        return None
    try:
        # Переводим запрос на английский через Claude
        translation_response = await anthropic.messages.create(
            model="claude-sonnet-4-5",
            max_tokens=200,
            messages=[{
                "role": "user",
                "content": f"Переведи этот математический запрос на английский язык для WolframAlpha. Верни ТОЛЬКО перевод без пояснений: {query}"
            }]
        )
        english_query = translation_response.content[0].text.strip()
        logger.info(f"WolframAlpha запрос: {english_query}")

        encoded = urllib.parse.quote(english_query)
        url = f"http://api.wolframalpha.com/v1/result?appid={WOLFRAM_API_KEY}&i={encoded}&units=metric"
        async with aiohttp.ClientSession() as s:
            async with s.get(url, timeout=aiohttp.ClientTimeout(total=15)) as resp:
                if resp.status == 200:
                    result = await resp.text()
                    return result
                else:
                    logger.warning(f"WolframAlpha статус: {resp.status}")
    except Exception as e:
        logger.error(f"WolframAlpha ошибка: {e}")
    return None


async def tavily_search(query: str) -> str | None:
    """Поиск актуальной информации через Tavily."""
    if not TAVILY_API_KEY:
        return None
    try:
        async with aiohttp.ClientSession() as s:
            async with s.post(
                "https://api.tavily.com/search",
                json={
                    "api_key": TAVILY_API_KEY,
                    "query": query,
                    "search_depth": "basic",
                    "max_results": 3
                },
                timeout=aiohttp.ClientTimeout(total=15)
            ) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    results = data.get("results", [])
                    if results:
                        text = ""
                        for r in results:
                            text += f"• {r.get('title', '')}: {r.get('content', '')[:300]}\n\n"
                        return text.strip()
    except Exception as e:
        logger.error(f"Tavily ошибка: {e}")
    return None


def is_search_request(text: str) -> bool:
    """Определяем запросы где нужна актуальная информация."""
    keywords = [
        "последние новости", "новости", "сейчас", "сегодня", "вчера",
        "в 2024", "в 2025", "в 2026", "актуально", "свежий", "недавно",
        "последний", "текущий", "нынешний", "на данный момент",
        "что происходит", "что случилось", "latest", "recent", "now", "today",
        "current", "news", "знаешь ли", "слышал ли", "в курсе",
        "нашли", "открыли", "обнаружили", "произошло", "случилось",
        "вышел", "вышла", "вышло", "появился", "появилась", "появилось",
        "объявили", "сообщили", "стало известно", "выяснилось",
        "расскажи про последние", "что нового", "какие новости"
    ]
    return any(kw in text.lower() for kw in keywords)


async def _generate_together(prompt: str) -> bytes | None:
    """Основной генератор — Together AI (FLUX)."""
    if not TOGETHER_API_KEY:
        return None
    url = "https://api.together.xyz/v1/images/generations"
    headers = {
        "Authorization": f"Bearer {TOGETHER_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": "black-forest-labs/FLUX.1-schnell-Free",
        "prompt": prompt,
        "width": 1024,
        "height": 1024,
        "n": 1,
        "response_format": "b64_json",
    }
    try:
        async with aiohttp.ClientSession() as s:
            async with s.post(url, headers=headers, json=payload,
                              timeout=aiohttp.ClientTimeout(total=60)) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    b64 = data["data"][0]["b64_json"]
                    return base64.b64decode(b64)
                else:
                    body = await resp.text()
                    logger.warning(f"Together: статус {resp.status}, {body[:200]}")
    except Exception as e:
        logger.error(f"Together ошибка: {type(e).__name__}: {e}")
    return None


async def _generate_pollinations(prompt: str) -> bytes | None:
    """Запасной генератор — Pollinations."""
    encoded = urllib.parse.quote(prompt, safe='')
    url = f"https://image.pollinations.ai/prompt/{encoded}?width=1024&height=1024&nologo=true"
    try:
        async with aiohttp.ClientSession() as s:
            async with s.get(url, timeout=aiohttp.ClientTimeout(total=90)) as resp:
                if resp.status == 200:
                    data = await resp.read()
                    content_type = resp.headers.get("Content-Type", "")
                    if data and content_type.startswith("image"):
                        return data
    except Exception as e:
        logger.error(f"Pollinations ошибка: {type(e).__name__}: {e}")
    return None


async def generate_image(prompt: str, retries: int = 2) -> bytes | None:
    """Пробует Together AI, при неудаче падает на Pollinations."""
    for attempt in range(1, retries + 1):
        # Сначала Together AI
        img = await _generate_together(prompt)
        if img:
            return img
        # Затем Pollinations
        img = await _generate_pollinations(prompt)
        if img:
            return img
        logger.warning(f"Генерация: попытка {attempt} не удалась")
        if attempt < retries:
            await asyncio.sleep(2)
    return None


def create_excel(text: str) -> bytes | None:
    try:
        lines = []
        in_table = False
        for line in text.split("\n"):
            if "EXCEL_TABLE:" in line:
                in_table = True
                continue
            if in_table and "|" in line:
                lines.append(line.strip())

        if not lines:
            return None

        wb = openpyxl.Workbook()
        ws = wb.active
        for i, line in enumerate(lines, 1):
            cells = line.split("|")
            for j, cell in enumerate(cells, 1):
                ws.cell(row=i, column=j, value=cell.strip())
                if i == 1:
                    ws.cell(row=i, column=j).font = openpyxl.styles.Font(bold=True)

        for col in ws.columns:
            max_len = max(len(str(cell.value or "")) for cell in col)
            ws.column_dimensions[col[0].column_letter].width = max_len + 4

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()
    except Exception as e:
        logger.error(f"Ошибка создания Excel: {e}")
        return None


def create_presentation(text: str) -> bytes | None:
    try:
        title = "Презентация"
        slides_data = []

        for line in text.split("\n"):
            line = line.strip()
            if line.startswith("TITLE:"):
                title = line[6:]
            elif line.startswith("SLIDE:"):
                parts = line[6:].split("|", 1)
                slide_title = parts[0] if parts else ""
                slide_content = parts[1] if len(parts) > 1 else ""
                slides_data.append((slide_title, slide_content))

        if not slides_data:
            return None

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Создано AI-ассистентом"

        for slide_title, slide_content in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_content

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except Exception as e:
        logger.error(f"Ошибка создания презентации: {e}")
        return None


async def ask_claude(user_id: int, content) -> str:
    history = await load_history(user_id)
    content_str = content if isinstance(content, str) else "[фото]"
    history.append({"role": "user", "content": content})

    response = await anthropic.messages.create(
        model="claude-sonnet-4-5",
        max_tokens=2048,
        system=await get_system_prompt(user_id),
        messages=history
    )

    reply = response.content[0].text

    await save_message(user_id, "user", content_str)
    await save_message(user_id, "assistant", reply)

    return reply


# ─── Хэндлеры ─────────────────────────────────────────────────────────────────
@dp.message(CommandStart())
async def cmd_start(message: Message):
    user_id = message.from_user.id
    username = message.from_user.username or message.from_user.first_name

    # Проверяем реферальную ссылку
    referred_by = None
    args = message.text.split()
    if len(args) > 1:
        try:
            referred_by = int(args[1])
            if referred_by == user_id:
                referred_by = None  # нельзя пригласить самого себя
        except ValueError:
            pass

    reward = await create_user(user_id, username, referred_by)

    # Активируем пробный период для новых пользователей
    trial_activated = await activate_trial(user_id)

    # Уведомляем реферера о новом приглашённом
    if referred_by:
        ref_count = await get_referral_count(referred_by)
        remaining = 5 - (ref_count % 5)
        try:
            if reward:
                await bot.send_message(
                    referred_by,
                    f"🎉 Ты пригласил уже {ref_count} друзей и получил бесплатный месяц подписки!\n"
                    f"Продолжай приглашать — каждые 5 друзей = 1 месяц бесплатно!"
                )
            else:
                await bot.send_message(
                    referred_by,
                    f"👥 По твоей ссылке зарегистрировался новый пользователь!\n"
                    f"Приглашено: {ref_count} из 5 для бесплатного месяца. Осталось: {remaining}!"
                )
        except Exception:
            pass

    sub_until = await get_subscription_until(user_id)
    if trial_activated:
        sub_text = f"🎁 Пробный период активирован на {TRIAL_DAYS} дня — безлимитный доступ!"
    elif sub_until:
        sub_text = f"✅ Подписка активна до {sub_until}"
    else:
        sub_text = f"🆓 Бесплатно: {FREE_MESSAGES_PER_DAY} сообщений/день"

    await message.answer(
        f"👋 Привет, {username}! Я AI-ассистент NeuroBot.\n\n"
        f"{sub_text}\n\n"
        "🔧 Я умею:\n"
        "• 💬 Отвечать на вопросы\n"
        "• 🖼 Распознавать картинки и решать задачи с фото\n"
        "• 🎨 Генерировать изображения по описанию\n"
        "• 📊 Создавать Excel таблицы\n"
        "• 📋 Создавать презентации PowerPoint\n"
        "• 📝 Создавать Word документы\n\n"
        "📢 Следи за обновлениями: @NeuroBot_Info\n\n"
        "👇 Используй меню ниже или просто напиши мне!",
        reply_markup=bottom_menu()
    )
    await message.answer("Открываю меню:", reply_markup=main_keyboard())


@dp.message(Command("givesubscription"))
async def cmd_give_subscription(message: Message):
    if message.from_user.id not in ADMIN_IDS:
        await message.answer("❌ У тебя нет доступа к этой команде.")
        return

    args = message.text.split()
    if len(args) < 2:
        await message.answer(
            "Использование: /givesubscription @username или /givesubscription user_id\n\n"
            "Пример: /givesubscription @swizyy1"
        )
        return

    target = args[1].replace("@", "")

    async with db_pool.acquire() as conn:
        # Ищем по username или user_id
        if target.isdigit():
            row = await conn.fetchrow("SELECT user_id, username FROM users WHERE user_id = $1", int(target))
        else:
            row = await conn.fetchrow("SELECT user_id, username FROM users WHERE username = $1", target)

    if not row:
        await message.answer(f"❌ Пользователь {target} не найден в базе.")
        return

    user_id = row["user_id"]
    username = row["username"] or str(user_id)

    await activate_subscription(user_id, from_payment=True)
    until = await get_subscription_until(user_id)

    await message.answer(f"✅ Подписка выдана @{username} до {until}!")

    try:
        await bot.send_message(
            user_id,
            f"🎉 Тебе выдана подписка NeuroBot!\n\n"
            f"✅ Подписка активна до {until}\n\n"
            f"Безлимитный доступ ко всем функциям!"
        )
    except Exception:
        pass


@dp.message(Command("referrals"))
async def cmd_referrals(message: Message):
    if message.from_user.id not in ADMIN_IDS:
        await message.answer("❌ У тебя нет доступа к этой команде.")
        return

    async with db_pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT username, referral_count FROM users
            WHERE referral_count > 0
            ORDER BY referral_count DESC
        """)

    if not rows:
        await message.answer("📊 Пока никто не привёл друзей.")
        return

    text = "📊 Рефералы:\n\n"
    for row in rows:
        username = row["username"] or "unknown"
        count = row["referral_count"]
        medal = "🏆" if count >= 10 else "🥈" if count >= 5 else "👤"
        text += f"{medal} @{username} — {count} друзей\n"

    await message.answer(text)


@dp.message(Command("referral"))
async def cmd_referral(message: Message):
    user_id = message.from_user.id
    ref_count = await get_referral_count(user_id)
    remaining = 5 - (ref_count % 5)
    bot_info = await bot.get_me()
    ref_link = f"https://t.me/{bot_info.username}?start={user_id}"

    await message.answer(
        f"👥 Реферальная программа\n\n"
        f"Приглашай друзей и получай бесплатные месяцы!\n\n"
        f"🎁 Каждые 5 приглашённых друзей = 1 месяц бесплатно\n\n"
        f"📊 Твоя статистика:\n"
        f"• Приглашено друзей: {ref_count}\n"
        f"• До следующего бесплатного месяца: {remaining}\n\n"
        f"🔗 Твоя ссылка:\n`{ref_link}`\n\n"
        f"Поделись ссылкой с друзьями — когда они зарегистрируются, ты получишь зачёт!",
        parse_mode="Markdown"
    )


# ─── Фоновая задача: уведомления об истечении подписки ───────────────────────
async def check_expiring_subscriptions():
    """Каждые 12 часов проверяем истекающие подписки."""
    while True:
        try:
            await reset_notified_expired()
            expiring = await get_users_expiring_soon()
            for user_id, username, sub_until in expiring:
                until_date = datetime.fromisoformat(sub_until).strftime("%d.%m.%Y")
                try:
                    await bot.send_message(
                        user_id,
                        f"⚠️ Твоя подписка истекает {until_date}!\n\n"
                        f"Продли подписку чтобы не потерять безлимитный доступ.",
                        reply_markup=subscription_keyboard()
                    )
                    await mark_notified(user_id)
                    logger.info(f"Уведомление отправлено пользователю {user_id}")
                except Exception as e:
                    logger.error(f"Ошибка отправки уведомления {user_id}: {e}")
        except Exception as e:
            logger.error(f"Ошибка проверки подписок: {e}")
        await asyncio.sleep(43200)  # 12 часов


async def daily_reset_messages():
    """Каждую ночь в 00:00 UTC сбрасываем счётчик сообщений."""
    while True:
        now = datetime.now()
        # Ждём до следующей полуночи UTC
        next_midnight = (now + timedelta(days=1)).replace(
            hour=0, minute=0, second=0, microsecond=0
        )
        wait_seconds = (next_midnight - now).total_seconds()
        await asyncio.sleep(wait_seconds)

        try:
            async with db_pool.acquire() as conn:
                today = datetime.now().date().isoformat()
                await conn.execute(
                    "UPDATE users SET messages_today = 0 WHERE last_message_date != $1",
                    today
                )
            logger.info("Ежедневный сброс счётчиков выполнен")
        except Exception as e:
            logger.error(f"Ошибка сброса счётчиков: {e}")
    """Каждые 12 часов проверяем истекающие подписки."""
    while True:
        try:
            await reset_notified_expired()
            expiring = await get_users_expiring_soon()
            for user_id, username, sub_until in expiring:
                until_date = datetime.fromisoformat(sub_until).strftime("%d.%m.%Y")
                try:
                    await bot.send_message(
                        user_id,
                        f"⚠️ Твоя подписка истекает {until_date}!\n\n"
                        f"Продли подписку чтобы не потерять безлимитный доступ.",
                        reply_markup=subscription_keyboard()
                    )
                    await mark_notified(user_id)
                    logger.info(f"Уведомление отправлено пользователю {user_id}")
                except Exception as e:
                    logger.error(f"Ошибка отправки уведомления {user_id}: {e}")
        except Exception as e:
            logger.error(f"Ошибка проверки подписок: {e}")
        await asyncio.sleep(43200)  # 12 часов


@dp.message(Command("clear"))
async def cmd_clear(message: Message):
    await clear_history(message.from_user.id)
    await message.answer("🗑️ История очищена!")


# ─── Обработчики кнопок меню ──────────────────────────────────────────────────
@dp.callback_query(F.data == "menu_subscribe")
async def menu_subscribe(callback):
    await callback.answer()
    sub_until = await get_subscription_until(callback.from_user.id)
    extra = f"\nТвоя подписка будет продлена от {sub_until}." if sub_until else ""
    await callback.message.answer(
        f"💳 Подписка на NeuroBot\n\n"
        f"• Безлимитные сообщения\n"
        f"• История диалогов сохраняется\n"
        f"• Генерация картинок без ограничений\n"
        f"• Решение задач с фото\n"
        f"• Все функции без ограничений\n\n"
        f"💰 Цена: 299₽ / месяц{extra}",
        reply_markup=subscription_keyboard()
    )


@dp.callback_query(F.data == "menu_status")
async def menu_status(callback):
    await callback.answer()
    user_id = callback.from_user.id
    sub_until = await get_subscription_until(user_id)
    if sub_until:
        await callback.message.answer(f"✅ Подписка активна до {sub_until}\n\nБез ограничений на сообщения!")
    else:
        can, remaining = await can_send_message(user_id)
        await callback.message.answer(
            f"🆓 Бесплатный план\n"
            f"Осталось сообщений сегодня: {remaining}/{FREE_MESSAGES_PER_DAY}\n\n"
            f"Купи подписку для безлимитного доступа!",
            reply_markup=subscription_keyboard()
        )


@dp.callback_query(F.data == "menu_referral")
async def menu_referral(callback):
    await callback.answer()
    user_id = callback.from_user.id
    ref_count = await get_referral_count(user_id)
    remaining = 5 - (ref_count % 5)
    bot_info = await bot.get_me()
    ref_link = f"https://t.me/{bot_info.username}?start={user_id}"
    await callback.message.answer(
        f"👥 Реферальная программа\n\n"
        f"🎁 Каждые 5 приглашённых друзей = 1 месяц бесплатно\n\n"
        f"📊 Твоя статистика:\n"
        f"• Приглашено друзей: {ref_count}\n"
        f"• До следующего бесплатного месяца: {remaining}\n\n"
        f"🔗 Твоя ссылка:\n`{ref_link}`",
        parse_mode="Markdown"
    )


@dp.callback_query(F.data == "menu_clear")
async def menu_clear(callback):
    await callback.answer()
    await clear_history(callback.from_user.id)
    await callback.message.answer("🗑️ История очищена!")


@dp.callback_query(F.data == "menu_speech")
async def menu_speech(callback):
    await callback.answer()
    user_id = callback.from_user.id
    current_style = await get_user_style(user_id)
    current_name = STYLE_NAMES.get(current_style, "😊 Как с другом")
    await callback.message.answer(
        f"🎙️ Настройки речи\n\n"
        f"Текущий стиль: {current_name}\n\n"
        f"Выбери как бот будет общаться с тобой:",
        reply_markup=speech_style_keyboard()
    )


@dp.callback_query(F.data.startswith("style_"))
async def set_style(callback):
    await callback.answer()
    style = callback.data.replace("style_", "")
    if style not in STYLE_PROMPTS:
        return
    await set_user_style(callback.from_user.id, style)
    style_name = STYLE_NAMES.get(style, style)
    await callback.message.answer(
        f"✅ Стиль общения изменён на {style_name}!\n\n"
        f"Теперь я буду общаться с тобой в этом стиле."
    )


@dp.message(Command("status"))
async def cmd_status(message: Message):
    user_id = message.from_user.id
    sub_until = await get_subscription_until(user_id)

    if sub_until:
        await message.answer(f"✅ Подписка активна до {sub_until}\n\nБез ограничений на сообщения!")
    else:
        can, remaining = await can_send_message(user_id)
        await message.answer(
            f"🆓 Бесплатный план\n"
            f"Осталось сообщений сегодня: {remaining}/{FREE_MESSAGES_PER_DAY}\n\n"
            f"Купи подписку для безлимитного доступа!",
            reply_markup=subscription_keyboard()
        )


@dp.message(Command("subscribe"))
async def cmd_subscribe(message: Message):
    sub_until = await get_subscription_until(message.from_user.id)
    extra = f"\nТвоя подписка будет продлена до следующего месяца от {sub_until}." if sub_until else ""

    await message.answer(
        f"💳 Подписка на NeuroBot\n\n"
        f"• Безлимитные сообщения\n"
        f"• История диалогов сохраняется\n"
        f"• Генерация картинок без ограничений\n"
        f"• Решение задач с фото\n"
        f"• Все функции без ограничений\n\n"
        f"💰 Цена: 299₽ / месяц{extra}",
        reply_markup=subscription_keyboard()
    )


@dp.message(Command("help"))
async def cmd_help(message: Message):
    await message.answer(
        "🤖 Примеры запросов:\n\n"
        "🎨 Генерация картинки:\n"
        "«Нарисуй закат над морем»\n\n"
        "📊 Excel таблица:\n"
        "«Создай таблицу с расходами за месяц»\n\n"
        "📋 Презентация:\n"
        "«Создай презентацию о космосе на 5 слайдов»\n\n"
        "📝 Word документ:\n"
        "«Напиши реферат о Второй мировой войне»\n"
        "«Создай договор аренды»\n"
        "«Напиши эссе на тему...»\n\n"
        "🖼 Анализ картинки:\n"
        "Просто отправь фото с подписью или без!\n\n"
        "📐 Математика и задачи:\n"
        "• Простые вычисления — точно ✅\n"
        "• Уравнения и формулы — точно ✅\n"
        "• Задачи с фото — хорошо, но проверяй сложные ⚠️\n"
        "• Задачи с чертежами — старается, но может ошибиться ⚠️\n\n"
        "💡 Если ответ неверный — напиши «неправильно» и пришли задачу заново!"
    )


@dp.message(Command("admin"))
async def cmd_admin(message: Message):
    if message.from_user.id not in ADMIN_IDS:
        await message.answer("❌ У тебя нет доступа к этой команде.")
        return

    async with db_pool.acquire() as conn:
        today = datetime.now().date().isoformat()
        now = datetime.now().isoformat()
        week_ago = (datetime.now() - timedelta(days=7)).isoformat()
        month_start = datetime.now().replace(day=1).date().isoformat()

        total_users = await conn.fetchval("SELECT COUNT(*) FROM users")
        active_today = await conn.fetchval("SELECT COUNT(*) FROM users WHERE last_message_date = $1", today)
        subscribers = await conn.fetchval("SELECT COUNT(*) FROM users WHERE subscription_until > $1", now)
        new_week = await conn.fetchval("SELECT COUNT(*) FROM users WHERE created_at > $1", week_ago)

        top_users = await conn.fetch(
            "SELECT username, messages_today FROM users ORDER BY messages_today DESC LIMIT 5"
        )
        sub_list = await conn.fetch(
            "SELECT username, subscription_until FROM users WHERE subscription_until > $1 ORDER BY subscription_until DESC",
            now
        )

        rev_today = await conn.fetchrow(
            "SELECT COUNT(*), SUM(amount) FROM payments WHERE created_at >= $1", today
        )
        rev_month = await conn.fetchrow(
            "SELECT COUNT(*), SUM(amount) FROM payments WHERE created_at >= $1", month_start
        )
        rev_total = await conn.fetchrow("SELECT COUNT(*), SUM(amount) FROM payments")

    revenue_today_count = rev_today[0] or 0
    revenue_today = rev_today[1] or 0
    revenue_month_count = rev_month[0] or 0
    revenue_month = rev_month[1] or 0
    revenue_total_count = rev_total[0] or 0
    revenue_total = rev_total[1] or 0

    top_text = "\n".join([f"  • {u['username'] or 'unknown'} — {u['messages_today']} сообщ." for u in top_users]) or "нет данных"
    sub_text = "\n".join([
        f"  • {u['username'] or 'unknown'} до {datetime.fromisoformat(u['subscription_until']).strftime('%d.%m.%Y')}"
        for u in sub_list
    ]) or "нет подписчиков"

    await message.answer(
        f"📊 Админ-панель NeuroBot\n"
        f"{'─' * 30}\n\n"
        f"👥 Всего пользователей: {total_users}\n"
        f"🟢 Активных сегодня: {active_today}\n"
        f"🆕 Новых за 7 дней: {new_week}\n"
        f"💎 Подписчиков: {subscribers}\n\n"
        f"💰 Выручка сегодня: {revenue_today}₽ ({revenue_today_count} платежей)\n"
        f"📅 Выручка за месяц: {revenue_month}₽ ({revenue_month_count} платежей)\n"
        f"🏦 Всего заработано: {revenue_total}₽ ({revenue_total_count} платежей)\n\n"
        f"🏆 Топ активных сегодня:\n{top_text}\n\n"
        f"💳 Подписчики:\n{sub_text}"
    )


# ─── Оплата через ЮKassa ──────────────────────────────────────────────────────
@dp.callback_query(F.data == "buy_subscription")
async def process_buy(callback):
    await callback.answer()
    user_id = callback.from_user.id
    username = callback.from_user.username or callback.from_user.first_name

    await callback.message.answer("⏳ Создаю ссылку для оплаты...")

    payment = await create_yukassa_payment(user_id, username)
    if payment:
        pending_payments[payment["payment_id"]] = user_id
        pay_keyboard = InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(
                text=f"💳 Оплатить {SUBSCRIPTION_PRICE_RUB}₽",
                url=payment["url"]
            )
        ]])
        await callback.message.answer(
            f"💳 Подписка NeuroBot — {SUBSCRIPTION_PRICE_RUB}₽/месяц\n\n"
            f"Нажми кнопку ниже для оплаты. После оплаты подписка активируется автоматически в течение минуты!",
            reply_markup=pay_keyboard
        )
        # Запускаем проверку платежа в фоне
        asyncio.create_task(poll_payment(payment["payment_id"], user_id))
    else:
        await callback.message.answer(
            "❌ Не удалось создать ссылку для оплаты. Попробуй позже или напиши администратору."
        )


# ─── Проверка лимитов ──────────────────────────────────────────────────────────
async def check_limits(message: Message) -> bool:
    """Возвращает True если пользователь может отправить сообщение"""
    if message.from_user.id in ADMIN_IDS:
        return True  # админ без лимитов
    user_id = message.from_user.id
    await create_user(user_id, message.from_user.username or message.from_user.first_name)

    # Проверяем ежедневный бонус
    if not await is_subscribed(user_id):
        bonus_given = await claim_daily_bonus(user_id)
        if bonus_given:
            await message.answer(
                f"🎁 Ежедневный бонус! +{DAILY_BONUS_MESSAGES} сообщения сегодня за активность!"
            )

    can, remaining = await can_send_message(user_id)

    if not can:
        await message.answer(
            f"⛔ Ты израсходовал лимит бесплатных сообщений на сегодня ({FREE_MESSAGES_PER_DAY} шт.)\n\n"
            f"Лимит обновится завтра, или купи подписку для безлимитного доступа!",
            reply_markup=subscription_keyboard()
        )
        return False

    if not await is_subscribed(user_id) and remaining <= 3:
        await message.answer(
            f"⚠️ Осталось {remaining} бесплатных сообщений на сегодня.",
            reply_markup=subscription_keyboard()
        )

    await increment_message_count(user_id)

    return True


# ─── Обработчики сообщений ────────────────────────────────────────────────────
@dp.message(F.photo)
async def handle_photo(message: Message):
    if not await check_limits(message):
        return

    await bot.send_chat_action(message.chat.id, "typing")
    try:
        photo = message.photo[-1]
        file = await bot.get_file(photo.file_id)
        buf = io.BytesIO()
        await bot.download_file(file.file_path, buf)
        photo_bytes = buf.getvalue()

        photo_b64 = base64.standard_b64encode(photo_bytes).decode()
        caption = message.caption or "Если на фото математическая или геометрическая задача — реши её строго пошагово, записывая все формулы и промежуточные вычисления. Проверь ответ в конце. Если это не задача — опиши что на изображении."

        content = [
            {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": photo_b64}},
            {"type": "text", "text": caption}
        ]

        reply = await ask_claude(message.from_user.id, content)
        await send_long_message(message, reply)

    except Exception as e:
        logger.error(f"Ошибка обработки фото: {e}")
        await message.answer("❌ Не удалось обработать изображение. Попробуй ещё раз.")


@dp.message(F.voice)
async def handle_voice(message: Message):
    if not await check_limits(message):
        return

    await bot.send_chat_action(message.chat.id, "typing")
    try:
        # Скачиваем голосовое сообщение
        file = await bot.get_file(message.voice.file_id)
        buf = io.BytesIO()
        await bot.download_file(file.file_path, buf)
        voice_bytes = buf.getvalue()
        voice_b64 = base64.standard_b64encode(voice_bytes).decode()

        # Транскрибируем через Claude
        transcribe_response = await anthropic.messages.create(
            model="claude-sonnet-4-5",
            max_tokens=500,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Это голосовое сообщение в формате ogg. Транскрибируй его текст. Верни ТОЛЬКО текст без пояснений."
                    }
                ]
            }]
        )
        # Поскольку Claude не слышит аудио напрямую, просим пользователя написать текст
        await message.answer(
            "🎤 Голосовые сообщения пока в разработке!\n\n"
            "Напиши свой вопрос текстом — отвечу быстро 😊"
        )

    except Exception as e:
        logger.error(f"Ошибка обработки голосового: {e}")
        await message.answer("❌ Не удалось обработать голосовое. Напиши текстом!")


@dp.message(Command("broadcast"))
async def cmd_broadcast(message: Message):
    if message.from_user.id not in ADMIN_IDS:
        await message.answer("❌ У тебя нет доступа к этой команде.")
        return

    text = message.text.replace("/broadcast", "").strip()
    if not text:
        await message.answer(
            "📢 Использование: /broadcast текст сообщения\n\n"
            "Пример: /broadcast 🎉 Добавили новую функцию!"
        )
        return

    await message.answer("📤 Начинаю рассылку...")

    async with db_pool.acquire() as conn:
        users = await conn.fetch("SELECT user_id FROM users")

    sent = 0
    failed = 0
    for row in users:
        try:
            await bot.send_message(row["user_id"], text)
            sent += 1
            await asyncio.sleep(0.05)
        except Exception:
            failed += 1

    await message.answer(
        f"✅ Рассылка завершена!\n\n"
        f"📨 Отправлено: {sent}\n"
        f"❌ Не доставлено: {failed}"
    )


async def send_long_message(message: Message, text: str):
    """Отправляет длинное сообщение разбивая на части по 4096 символов."""
    formatted = format_response(text)
    max_len = 4096
    if len(formatted) <= max_len:
        try:
            await message.answer(formatted, parse_mode="Markdown")
        except Exception:
            await message.answer(formatted)
        return
    parts = [formatted[i:i+max_len] for i in range(0, len(formatted), max_len)]
    for part in parts:
        try:
            await message.answer(part, parse_mode="Markdown")
        except Exception:
            await message.answer(part)


def format_response(text: str) -> str:
    """
    Убирает лишний markdown из обычного текста,
    но сохраняет блоки кода в красивых плашках.
    """
    import re

    result = []
    # Разбиваем текст на части: код и не-код
    parts = re.split(r"(```[\w]*\n?.*?```)", text, flags=re.DOTALL)

    for part in parts:
        if part.startswith("```"):
            # Это блок кода — оставляем как есть
            result.append(part)
        else:
            # Обычный текст — убираем markdown
            part = re.sub(r"\*\*(.+?)\*\*", r"\1", part)
            part = re.sub(r"\*(.+?)\*", r"\1", part)
            part = re.sub(r"__(.+?)__", r"\1", part)
            part = re.sub(r"^#{1,6}\s+", "", part, flags=re.MULTILINE)
            result.append(part)

    return "".join(result).strip()


@dp.message(F.text)
async def handle_text(message: Message):
    if not await check_limits(message):
        return

    user_id = message.from_user.id
    text = message.text

    await bot.send_chat_action(message.chat.id, "typing")

    try:
        # Обработка кнопки меню
        if text == "📋 Меню":
            await message.answer("Открываю меню:", reply_markup=main_keyboard())
            return

        # Определяем что пользователь недоволен ответом
        retry_keywords = [
            "неправильно", "неверно", "ошибка", "ошибся", "не так",
            "перерешай", "попробуй снова", "попробуй ещё", "реши заново",
            "wrong", "incorrect", "try again", "redo", "mistake"
        ]
        if any(kw in text.lower() for kw in retry_keywords):
            clear_history(user_id)
            await message.answer(
                "🔄 Понял, давай попробуем заново с чистого листа!\n\n"
                "Пришли задачу ещё раз — решу пошагово и аккуратно."
            )
            return
        if is_image_request(text):
            await message.answer("🎨 Генерирую изображение, подожди 10-20 секунд...")
            await bot.send_chat_action(message.chat.id, "upload_photo")
            image_bytes = await generate_image(text)
            if image_bytes:
                await message.answer_photo(
                    BufferedInputFile(image_bytes, filename="image.jpg"),
                    caption="✅ Готово!"
                )
            else:
                await message.answer("❌ Не удалось сгенерировать изображение. Попробуй ещё раз.")
            return

        if is_excel_request(text) or is_presentation_request(text) or is_word_request(text):
            reply = await ask_claude(user_id, text)

            if "EXCEL_TABLE:" in reply:
                excel_bytes = create_excel(reply)
                if excel_bytes:
                    await message.answer_document(
                        BufferedInputFile(excel_bytes, filename="таблица.xlsx"),
                        caption="📊 Вот твоя Excel таблица!"
                    )
                    return

            if "PRESENTATION:" in reply:
                pptx_bytes = create_presentation(reply)
                if pptx_bytes:
                    await message.answer_document(
                        BufferedInputFile(pptx_bytes, filename="презентация.pptx"),
                        caption="📋 Вот твоя презентация!"
                    )
                    return

            if is_word_request(text):
                word_bytes = create_word(reply)
                if word_bytes:
                    await message.answer_document(
                        BufferedInputFile(word_bytes, filename="документ.docx"),
                        caption="📝 Вот твой Word документ!"
                    )
                    return

            await send_long_message(message, reply)
            return

        # Пробуем WolframAlpha для точных вычислений
        if is_math_request(text) and WOLFRAM_API_KEY:
            wolfram_result = await wolfram_calculate(text)
            if wolfram_result:
                enhanced_text = (
                    f"{text}\n\n"
                    f"[Точный ответ от WolframAlpha: {wolfram_result}]\n"
                    f"Объясни решение пошагово, используя этот точный ответ."
                )
                reply = await ask_claude(user_id, enhanced_text)
                await send_long_message(message, f"🔢 Точный ответ: {wolfram_result}\n\n{reply}")
                return

        # Поиск актуальной информации через Tavily
        if is_search_request(text) and TAVILY_API_KEY:
            search_result = await tavily_search(text)
            if search_result:
                enhanced_text = (
                    f"{text}\n\n"
                    f"[Актуальная информация из интернета:\n{search_result}]\n"
                    f"Ответь на вопрос используя эту актуальную информацию."
                )
                reply = await ask_claude(user_id, enhanced_text)
                await send_long_message(message, reply)
                return

        reply = await ask_claude(user_id, text)
        await send_long_message(message, reply)

    except Exception as e:
        logger.error(f"Ошибка: {e}")
        await message.answer("❌ Произошла ошибка. Попробуй ещё раз.")


# ─── Запуск ───────────────────────────────────────────────────────────────────
async def main():
    await init_db()
    threading.Thread(target=run_health_server, daemon=True).start()
    logger.info("Бот запускается...")
    asyncio.create_task(check_expiring_subscriptions())
    asyncio.create_task(daily_reset_messages())
    await dp.start_polling(bot)


if __name__ == "__main__":
    asyncio.run(main())
